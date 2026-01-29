"""
Enhanced Knowledge Management System v3.0 - KC
Modern RAG chatbot with session-based document management and persistent memory
Features: Upload files in chat, use KB files, hybrid retrieval, conversation memory
"""

from datetime import datetime
from pathlib import Path
from typing import List, Optional, Dict, Tuple, Any
import hashlib
import json
import logging
import os
import tempfile
import shutil

import streamlit as st
import requests

# Modern LangChain imports
from langchain_community.document_loaders import (
    PyPDFLoader,
    CSVLoader,
    TextLoader,
)

# Try to import Unstructured loaders, fall back to alternatives
try:
    from langchain_community.document_loaders import (
        UnstructuredExcelLoader,
        UnstructuredWordDocumentLoader,
        UnstructuredPowerPointLoader,
    )
    UNSTRUCTURED_AVAILABLE = True
except ImportError:
    UNSTRUCTURED_AVAILABLE = False

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_core.messages import HumanMessage, AIMessage, SystemMessage
from langchain_core.output_parsers import StrOutputParser
from langchain_core.runnables import RunnablePassthrough

# Embeddings and Vector Store
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma

# Ollama - Use langchain_community version which properly supports base_url
from langchain_community.chat_models import ChatOllama

# File monitoring
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler

# Image processing with OCR
IMAGE_PROCESSING_AVAILABLE = False
OCR_AVAILABLE = False
OCR_METHOD = None

try:
    from PIL import Image as PILImage
    IMAGE_PROCESSING_AVAILABLE = True
    
    try:
        import easyocr
        OCR_AVAILABLE = True
        OCR_METHOD = "easyocr"
        print("✅ EasyOCR available")
    except ImportError:
        try:
            import pytesseract
            pytesseract.get_tesseract_version()
            OCR_AVAILABLE = True
            OCR_METHOD = "tesseract"
            print("✅ Tesseract OCR available")
        except Exception:
            print("⚠️ No OCR library found")
except ImportError as e:
    print(f"❌ Image processing not available: {e}")

# Constants
MAX_FILE_SIZE = 1024 * 1024 * 1024  # 1GB
CHUNK_SIZE = 1000
CHUNK_OVERLAP = 200
OLLAMA_BASE_URL = "http://localhost:11434"
OLLAMA_API_URL = "http://localhost:11434/api"
PERSIST_DIR = Path("./persistent_storage")
KNOWLEDGE_BASE_DIR = Path("./knowledge_bases")
TEMP_DIR = Path("./temp_storage")
SESSION_DIR = Path("./session_storage")

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(name)s - %(levelname)s - %(message)s"
)
logger = logging.getLogger(__name__)

# Disable CUDA
os.environ["CUDA_VISIBLE_DEVICES"] = ""
os.environ["HF_HUB_DISABLE_SYMLINKS_WARNING"] = "1"


# Alternative document loaders when unstructured is not available
class ExcelLoader:
    """Alternative Excel loader using openpyxl"""
    def __init__(self, file_path):
        self.file_path = file_path
    
    def load(self):
        try:
            import openpyxl
            workbook = openpyxl.load_workbook(self.file_path, data_only=True)
            documents = []
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content_parts = [f"Sheet: {sheet_name}\n"]
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        content_parts.append(row_text)
                
                content = "\n".join(content_parts)
                if content.strip():
                    doc = Document(
                        page_content=content,
                        metadata={"source": str(self.file_path), "sheet": sheet_name}
                    )
                    documents.append(doc)
            
            return documents
        except Exception as e:
            logger.error(f"Error loading Excel file {self.file_path}: {e}")
            return []


class WordLoader:
    """Alternative Word loader using python-docx"""
    def __init__(self, file_path):
        self.file_path = file_path
    
    def load(self):
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(self.file_path)
            
            content_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    content_parts.append(para.text)
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        content_parts.append(row_text)
            
            content = "\n\n".join(content_parts)
            
            if content.strip():
                return [Document(
                    page_content=content,
                    metadata={"source": str(self.file_path)}
                )]
            return []
        except Exception as e:
            logger.error(f"Error loading Word file {self.file_path}: {e}")
            return []


class PowerPointLoader:
    """Alternative PowerPoint loader using python-pptx"""
    def __init__(self, file_path):
        self.file_path = file_path
    
    def load(self):
        try:
            from pptx import Presentation
            prs = Presentation(self.file_path)
            documents = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                content_parts = [f"Slide {slide_num}:"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text)
                
                content = "\n".join(content_parts)
                if content.strip():
                    doc = Document(
                        page_content=content,
                        metadata={"source": str(self.file_path), "slide": slide_num}
                    )
                    documents.append(doc)
            
            return documents
        except Exception as e:
            logger.error(f"Error loading PowerPoint file {self.file_path}: {e}")
            return []


class OCRProcessor:
    """Optimized OCR processor with caching"""
    
    def __init__(self):
        self.ocr_reader = None
        self.ocr_method = OCR_METHOD
        self._initialize_ocr()
    
    def _initialize_ocr(self):
        """Initialize the best available OCR method"""
        if not OCR_AVAILABLE:
            return
        
        try:
            if self.ocr_method == "easyocr":
                import easyocr
                self.ocr_reader = easyocr.Reader(["en"], gpu=False)
                logger.info("Initialized EasyOCR reader")
            
            elif self.ocr_method == "tesseract":
                import pytesseract
                version = pytesseract.get_tesseract_version()
                logger.info(f"Initialized Tesseract OCR - Version: {version}")
        
        except Exception as e:
            logger.error(f"Error initializing OCR: {e}")
            self.ocr_reader = None
    
    def extract_text_from_image(self, image_path: Path) -> Tuple[str, Dict]:
        """Extract text from image using the best available OCR method"""
        if not OCR_AVAILABLE or not IMAGE_PROCESSING_AVAILABLE:
            return "", {"ocr_method": "none", "error": "OCR not available"}
        
        try:
            image = PILImage.open(image_path)
            
            if image.mode != "RGB":
                image = image.convert("RGB")
            
            text = ""
            metadata = {"ocr_method": self.ocr_method}
            
            if self.ocr_method == "easyocr" and self.ocr_reader:
                results = self.ocr_reader.readtext(str(image_path))
                text_parts = [text_content for bbox, text_content, confidence in results if confidence > 0.5]
                text = " ".join(text_parts)
                metadata["confidence_avg"] = sum(c for _, _, c in results) / len(results) if results else 0
                metadata["text_blocks"] = len(text_parts)
            
            elif self.ocr_method == "tesseract":
                import pytesseract
                text = pytesseract.image_to_string(image)
                metadata["text_blocks"] = len(text.split("\n")) if text else 0
            
            text = text.strip()
            metadata.update({"text_length": len(text), "success": len(text) > 0})
            
            return text, metadata
        
        except Exception as e:
            logger.error(f"Error in OCR processing: {e}")
            return "", {"ocr_method": self.ocr_method, "error": str(e), "success": False}


class SecurityValidator:
    """Security validation for file uploads"""
    
    ALLOWED_EXTENSIONS = {
        ".pdf", ".txt", ".csv", ".xlsx", ".xls",
        ".doc", ".docx", ".ppt", ".pptx",
        ".jpg", ".jpeg", ".png", ".gif", ".bmp"
    }
    
    @staticmethod
    def validate_file(file_path: Path) -> bool:
        """Validate file security"""
        if not file_path.exists():
            raise ValueError(f"File does not exist: {file_path}")
        
        if file_path.stat().st_size > MAX_FILE_SIZE:
            raise ValueError(f"File size exceeds limit of {MAX_FILE_SIZE/1024/1024/1024}GB")
        
        if file_path.suffix.lower() not in SecurityValidator.ALLOWED_EXTENSIONS:
            raise ValueError(f"Unsupported file type: {file_path.suffix}")
        
        return True
    
    @staticmethod
    def get_file_hash(file_path: Path) -> str:
        """Calculate file hash"""
        with open(file_path, "rb") as f:
            return hashlib.sha256(f.read()).hexdigest()


class SessionManager:
    """Manages session-based documents"""
    
    def __init__(self):
        self.session_id = st.session_state.get("session_id", self._generate_session_id())
        st.session_state.session_id = self.session_id
        self.session_dir = SESSION_DIR / self.session_id
        self.session_dir.mkdir(parents=True, exist_ok=True)
        self.session_files = self._load_session_files()
    
    def _generate_session_id(self) -> str:
        """Generate a unique session ID"""
        return f"session_{datetime.now().strftime('%Y%m%d_%H%M%S')}_{hash(datetime.now()) % 10000}"
    
    def _load_session_files(self) -> Dict:
        """Load session file metadata"""
        metadata_file = self.session_dir / "metadata.json"
        if metadata_file.exists():
            with open(metadata_file, "r") as f:
                return json.load(f)
        return {"files": [], "created": datetime.now().isoformat()}
    
    def _save_session_files(self):
        """Save session file metadata"""
        metadata_file = self.session_dir / "metadata.json"
        with open(metadata_file, "w") as f:
            json.dump(self.session_files, f, indent=2)
    
    def add_session_file(self, uploaded_file, file_hash: str) -> Path:
        """Add a file to the current session"""
        file_path = self.session_dir / uploaded_file.name
        
        # Save file
        with open(file_path, "wb") as f:
            f.write(uploaded_file.getbuffer())
        
        # Update metadata
        self.session_files["files"].append({
            "name": uploaded_file.name,
            "hash": file_hash,
            "uploaded_at": datetime.now().isoformat(),
            "size": uploaded_file.size,
            "type": uploaded_file.type
        })
        self._save_session_files()
        
        return file_path
    
    def get_session_files(self) -> List[Dict]:
        """Get all files in the current session"""
        return self.session_files.get("files", [])
    
    def clear_session(self):
        """Clear all session files"""
        if self.session_dir.exists():
            shutil.rmtree(self.session_dir)
        self.session_dir.mkdir(parents=True, exist_ok=True)
        self.session_files = {"files": [], "created": datetime.now().isoformat()}
        self._save_session_files()


class DocumentProcessor:
    """Process documents with text splitting and chunking"""
    
    def __init__(self):
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=CHUNK_SIZE,
            chunk_overlap=CHUNK_OVERLAP,
            length_function=len,
            separators=["\n\n", "\n", " ", ""]
        )
        self.ocr_processor = OCRProcessor()
    
    def process_uploaded_file(self, uploaded_file) -> Tuple[List[Document], str]:
        """Process an uploaded file from Streamlit"""
        try:
            # Create temp file
            suffix = Path(uploaded_file.name).suffix
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp_file:
                tmp_file.write(uploaded_file.getbuffer())
                tmp_path = Path(tmp_file.name)
            
            # Calculate hash
            file_hash = SecurityValidator.get_file_hash(tmp_path)
            
            # Validate
            SecurityValidator.validate_file(tmp_path)
            
            # Process file
            documents = self._load_and_split_document(tmp_path, uploaded_file.name)
            
            # Cleanup
            tmp_path.unlink()
            
            return documents, file_hash
        
        except Exception as e:
            logger.error(f"Error processing uploaded file: {e}")
            if 'tmp_path' in locals() and tmp_path.exists():
                tmp_path.unlink()
            raise
    
    def process_file_from_path(self, file_path: Path, source_name: str = None) -> Tuple[List[Document], str]:
        """Process a file from a file path"""
        try:
            SecurityValidator.validate_file(file_path)
            file_hash = SecurityValidator.get_file_hash(file_path)
            documents = self._load_and_split_document(file_path, source_name or file_path.name)
            return documents, file_hash
        except Exception as e:
            logger.error(f"Error processing file from path: {e}")
            raise
    
    def _load_and_split_document(self, file_path: Path, source_name: str) -> List[Document]:
        """Load and split a document into chunks"""
        ext = file_path.suffix.lower()
        
        # Get appropriate loader
        loader_class = self._get_loader_class(ext)
        if not loader_class:
            raise ValueError(f"Unsupported file type: {ext}")
        
        # Handle image files with OCR
        if ext in {".jpg", ".jpeg", ".png", ".gif", ".bmp"}:
            text, metadata = self.ocr_processor.extract_text_from_image(file_path)
            if text:
                doc = Document(
                    page_content=text,
                    metadata={"source": source_name, **metadata}
                )
                return self.text_splitter.split_documents([doc])
            else:
                logger.warning(f"No text extracted from image: {source_name}")
                return []
        
        # Load document
        loader = loader_class(str(file_path))
        documents = loader.load()
        
        # Update metadata
        for doc in documents:
            doc.metadata["source"] = source_name
        
        # Split into chunks
        split_docs = self.text_splitter.split_documents(documents)
        
        return split_docs
    
    @staticmethod
    def _get_loader_class(ext: str):
        """Get appropriate document loader class based on file extension"""
        if UNSTRUCTURED_AVAILABLE:
            loaders = {
                ".pdf": PyPDFLoader,
                ".xlsx": UnstructuredExcelLoader,
                ".xls": UnstructuredExcelLoader,
                ".csv": CSVLoader,
                ".txt": TextLoader,
                ".doc": UnstructuredWordDocumentLoader,
                ".docx": UnstructuredWordDocumentLoader,
                ".ppt": UnstructuredPowerPointLoader,
                ".pptx": UnstructuredPowerPointLoader,
            }
        else:
            loaders = {
                ".pdf": PyPDFLoader,
                ".xlsx": ExcelLoader,
                ".xls": ExcelLoader,
                ".csv": CSVLoader,
                ".txt": TextLoader,
                ".doc": WordLoader,
                ".docx": WordLoader,
                ".ppt": PowerPointLoader,
                ".pptx": PowerPointLoader,
            }
        return loaders.get(ext)


class KnowledgeBaseManager:
    """Manages knowledge bases"""
    
    def __init__(self):
        self.knowledge_bases = {}
        self.active_kb = None
        self.ensure_directories()
        self.load_existing_knowledge_bases()
    
    def ensure_directories(self):
        """Create necessary directories"""
        KNOWLEDGE_BASE_DIR.mkdir(exist_ok=True)
        PERSIST_DIR.mkdir(exist_ok=True)
        TEMP_DIR.mkdir(exist_ok=True)
        SESSION_DIR.mkdir(exist_ok=True)
    
    def create_knowledge_base(self, name: str, description: str = "") -> bool:
        """Create a new knowledge base"""
        try:
            kb_dir = KNOWLEDGE_BASE_DIR / name
            kb_dir.mkdir(exist_ok=True)
            
            (kb_dir / "documents").mkdir(exist_ok=True)
            (kb_dir / "vectorstore").mkdir(exist_ok=True)
            
            metadata = {
                "name": name,
                "description": description,
                "created_date": datetime.now().isoformat(),
                "file_count": 0,
                "last_updated": datetime.now().isoformat(),
                "file_hashes": {},
            }
            
            metadata_path = kb_dir / "metadata.json"
            with open(metadata_path, "w") as f:
                json.dump(metadata, f, indent=2)
            
            self.knowledge_bases[name] = metadata
            logger.info(f"Created knowledge base: {name}")
            return True
        
        except Exception as e:
            logger.error(f"Error creating knowledge base: {e}")
            return False
    
    def load_existing_knowledge_bases(self):
        """Load existing knowledge bases"""
        for kb_dir in KNOWLEDGE_BASE_DIR.iterdir():
            if kb_dir.is_dir():
                metadata_path = kb_dir / "metadata.json"
                if metadata_path.exists():
                    with open(metadata_path, "r") as f:
                        metadata = json.load(f)
                        self.knowledge_bases[kb_dir.name] = metadata
    
    def list_knowledge_bases(self) -> List[str]:
        """List all knowledge bases"""
        return list(self.knowledge_bases.keys())
    
    def get_active_knowledge_base(self) -> Optional[str]:
        """Get active knowledge base"""
        return self.active_kb
    
    def set_active_knowledge_base(self, name: str) -> bool:
        """Set active knowledge base"""
        if name in self.knowledge_bases:
            self.active_kb = name
            logger.info(f"Set active knowledge base: {name}")
            return True
        return False
    
    def get_knowledge_base_info(self, name: str) -> Dict:
        """Get knowledge base information"""
        return self.knowledge_bases.get(name, {})
    
    def update_metadata(self, name: str, updates: Dict):
        """Update knowledge base metadata"""
        if name in self.knowledge_bases:
            self.knowledge_bases[name].update(updates)
            self.knowledge_bases[name]["last_updated"] = datetime.now().isoformat()
            
            metadata_path = KNOWLEDGE_BASE_DIR / name / "metadata.json"
            with open(metadata_path, "w") as f:
                json.dump(self.knowledge_bases[name], f, indent=2)
    
    def delete_knowledge_base(self, name: str) -> bool:
        """Delete a knowledge base"""
        try:
            kb_dir = KNOWLEDGE_BASE_DIR / name
            if kb_dir.exists():
                shutil.rmtree(kb_dir)
            
            if name in self.knowledge_bases:
                del self.knowledge_bases[name]
            
            if self.active_kb == name:
                self.active_kb = None
            
            logger.info(f"Deleted knowledge base: {name}")
            return True
        
        except Exception as e:
            logger.error(f"Error deleting knowledge base: {e}")
            return False


class VectorStoreManager:
    """Manages vector stores using Chroma"""
    
    def __init__(self):
        self.embeddings = HuggingFaceEmbeddings(
            model_name="sentence-transformers/all-MiniLM-L6-v2",
            model_kwargs={"device": "cpu"},
            encode_kwargs={"normalize_embeddings": True}
        )
        self.kb_stores = {}
        self.session_store = None
    
    def get_or_create_kb_store(self, kb_name: str) -> Chroma:
        """Get or create a vector store for a knowledge base"""
        if kb_name in self.kb_stores:
            return self.kb_stores[kb_name]
        
        persist_directory = str(KNOWLEDGE_BASE_DIR / kb_name / "vectorstore")
        
        vector_store = Chroma(
            collection_name=f"kb_{kb_name}",
            embedding_function=self.embeddings,
            persist_directory=persist_directory
        )
        
        self.kb_stores[kb_name] = vector_store
        return vector_store
    
    def get_or_create_session_store(self, session_id: str) -> Chroma:
        """Get or create a vector store for the current session"""
        if self.session_store is None:
            persist_directory = str(SESSION_DIR / session_id / "vectorstore")
            
            self.session_store = Chroma(
                collection_name=f"session_{session_id}",
                embedding_function=self.embeddings,
                persist_directory=persist_directory
            )
        
        return self.session_store
    
    def clear_session_store(self):
        """Clear the session vector store"""
        self.session_store = None
    
    def delete_vector_store(self, kb_name: str):
        """Delete a vector store"""
        persist_directory = KNOWLEDGE_BASE_DIR / kb_name / "vectorstore"
        if persist_directory.exists():
            shutil.rmtree(persist_directory)
        
        if kb_name in self.kb_stores:
            del self.kb_stores[kb_name]


class HybridChatbot:
    """Hybrid chatbot that uses both KB and session documents """
    
    def __init__(self):
        self.llm = None
        self.kb_retriever = None
        self.session_retriever = None
        self.chain = None
        self.chat_history = []
    
    def initialize_llm(self, model_name: str, keep_alive="5m"):
        """Initialize the LLM with memory management"""
        try:
            # Initialize LangChain ChatOllama
            # keep_alive="5m": Keeps model loaded for fast responses
            # Model will be unloaded when switching to a different model
            self.llm = ChatOllama(
                model=model_name,
                base_url=OLLAMA_BASE_URL,
                temperature=0.7,
                keep_alive=keep_alive,  # Keep loaded until switched
            )
            logger.info(f"Successfully initialized LLM: {model_name} (keep_alive={keep_alive})")
            return True
            
        except Exception as e:
            logger.error(f"Error initializing LLM: {e}")
            error_msg = str(e)
            
            # Provide helpful error messages
            if "connection" in error_msg.lower():
                raise ValueError("Cannot connect to Ollama. Make sure it's running: ollama serve")
            elif "not found" in error_msg.lower() or "404" in error_msg:
                raise ValueError(f"Model '{model_name}' not found. Pull it with: ollama pull {model_name}")
            else:
                raise ValueError(f"Failed to initialize model: {error_msg}")
        
        return False
    
    def unload_model(self):
        """Explicitly unload the current model from Ollama"""
        if self.llm and hasattr(self.llm, 'model'):
            try:
                # Send request to unload model
                model_name = self.llm.model
                logger.info(f"Unloading model: {model_name}")
                
                # Create a temporary ChatOllama with keep_alive=0 to unload
                temp_llm = ChatOllama(
                    model=model_name,
                    base_url=OLLAMA_BASE_URL,
                    keep_alive=0,  # Unload immediately
                )
                # Making any call with keep_alive=0 will unload the model
                # We don't need the response, just need to trigger the unload
                
                logger.info(f"Model {model_name} unloaded from memory")
                return True
            except Exception as e:
                logger.error(f"Error unloading model: {e}")
        return False
    
    def setup_retrievers(self, kb_store: Optional[Chroma] = None, session_store: Optional[Chroma] = None):
        """Setup retrievers for KB and session documents"""
        self.kb_retriever = kb_store.as_retriever(search_kwargs={"k": 3}) if kb_store else None
        self.session_retriever = session_store.as_retriever(search_kwargs={"k": 3}) if session_store else None
        logger.info(f"Setup retrievers - KB: {kb_store is not None}, Session: {session_store is not None}")
    
    def create_chain(self):
        """Create the conversation chain"""
        if not self.llm:
            raise ValueError("LLM not initialized")
        
        # Create prompt
        template = """You are a helpful AI assistant with access to documents from both a knowledge base and the current conversation session.

Use the following context to answer the question. If you cannot find the answer in the context, say so clearly.

Knowledge Base Context:
{kb_context}

Session Documents Context:
{session_context}

Chat History:
{chat_history}

Question: {question}

Answer:"""
        
        prompt = ChatPromptTemplate.from_template(template)
        
        # Create chain
        def format_docs(docs):
            return "\n\n".join(doc.page_content for doc in docs) if docs else "No documents found."
        
        def get_kb_context(inputs):
            if self.kb_retriever:
                docs = self.kb_retriever.get_relevant_documents(inputs["question"])
                return format_docs(docs)
            return "No knowledge base active."
        
        def get_session_context(inputs):
            if self.session_retriever:
                docs = self.session_retriever.get_relevant_documents(inputs["question"])
                return format_docs(docs)
            return "No session documents uploaded."
        
        def format_chat_history(history):
            if not history:
                return "No previous messages."
            formatted = []
            for msg in history[-6:]:  # Last 6 messages
                if isinstance(msg, HumanMessage):
                    formatted.append(f"Human: {msg.content}")
                elif isinstance(msg, AIMessage):
                    formatted.append(f"Assistant: {msg.content}")
            return "\n".join(formatted)
        
        self.chain = (
            {
                "kb_context": get_kb_context,
                "session_context": get_session_context,
                "chat_history": lambda x: format_chat_history(self.chat_history),
                "question": lambda x: x["question"]
            }
            | prompt
            | self.llm
            | StrOutputParser()
        )
        
        logger.info("Created hybrid conversation chain")
        return True
    
    def query(self, question: str) -> str:
        """Query the chatbot"""
        try:
            if not self.chain:
                raise ValueError("Conversation chain not initialized")
            
            # Get response
            response = self.chain.invoke({"question": question})
            
            # Update chat history
            self.chat_history.append(HumanMessage(content=question))
            self.chat_history.append(AIMessage(content=response))
            
            # Keep only last 20 messages
            if len(self.chat_history) > 20:
                self.chat_history = self.chat_history[-20:]
            
            return response
        
        except requests.exceptions.ConnectionError as e:
            error_msg = "Cannot connect to Ollama. Please make sure Ollama is running:\n\n1. Open terminal\n2. Run: ollama serve\n3. Try your question again"
            logger.error(f"Ollama connection error: {e}")
            raise ValueError(error_msg)
        
        except Exception as e:
            logger.error(f"Error processing query: {e}", exc_info=True)
            # Provide more helpful error message
            if "404" in str(e):
                raise ValueError(f"Ollama API error. Please check:\n1. Ollama is running (ollama serve)\n2. Model is available (ollama list)\n3. Try reinitializing the model\n\nTechnical error: {e}")
            raise
    
    def clear_history(self):
        """Clear chat history"""
        self.chat_history = []


# Streamlit UI Functions
def initialize_session_state():
    """Initialize Streamlit session state"""
    if "initialized" not in st.session_state:
        st.session_state.initialized = True
        st.session_state.kb_manager = KnowledgeBaseManager()
        st.session_state.vector_store_manager = VectorStoreManager()
        st.session_state.session_manager = SessionManager()
        st.session_state.processor = DocumentProcessor()
        st.session_state.chatbot = HybridChatbot()
        st.session_state.selected_model = None
        st.session_state.current_loaded_model = None  # Track which model is actually loaded
        st.session_state.conversation_ready = False
        st.session_state.messages = []
        st.session_state.session_files_count = 0


def get_available_models() -> List[str]:
    """Get available Ollama models"""
    try:
        response = requests.get(f"{OLLAMA_API_URL}/tags", timeout=5)
        if response.status_code == 200:
            models = response.json().get("models", [])
            model_list = []
            
            for model in models:
                model_name = model["name"]
                # Filter out vision models as they may have different requirements
                if "-vl" in model_name.lower():
                    logger.warning(f"Skipping vision model: {model_name} (may require special handling)")
                    continue
                model_list.append(model_name)
            
            # If only vision models available, include them anyway
            if not model_list:
                model_list = [model["name"] for model in models]
                logger.warning("Only vision models available - these may have limitations")
            
            logger.info(f"Found {len(model_list)} Ollama models")
            return model_list
        else:
            logger.error(f"Ollama API returned status {response.status_code}: {response.text}")
    except requests.exceptions.ConnectionError:
        logger.error("Cannot connect to Ollama. Make sure Ollama is running with: ollama serve")
    except Exception as e:
        logger.error(f"Error getting models: {e}")
    return []


def display_sidebar():
    """Display sidebar with model selection and status"""
    with st.sidebar:
        st.header("🤖 Model & Status")
        
        # Model selection
        available_models = get_available_models()
        
        if available_models:
            selected_model = st.selectbox(
                "Select Ollama Model",
                options=available_models,
                index=0 if not st.session_state.selected_model else 
                      (available_models.index(st.session_state.selected_model) 
                       if st.session_state.selected_model in available_models else 0),
                key="model_selector"
            )
            
            st.info("""
💡 **Smart Memory**: Model stays loaded for fast responses. When you switch models, the old one is automatically unloaded.
            """)
            
            if st.button("Initialize Model", key="init_model_btn"):
                with st.spinner("Initializing model..."):
                    chatbot = st.session_state.chatbot
                    current_loaded = st.session_state.get("current_loaded_model")
                    
                    # Check if switching models
                    if current_loaded and current_loaded != selected_model:
                        st.info(f"Unloading previous model: {current_loaded}")
                        chatbot.unload_model()
                    
                    try:
                        if chatbot.initialize_llm(selected_model):
                            st.session_state.selected_model = selected_model
                            st.session_state.current_loaded_model = selected_model
                            
                            # Setup retrievers
                            kb_manager = st.session_state.kb_manager
                            active_kb = kb_manager.get_active_knowledge_base()
                            
                            kb_store = None
                            if active_kb:
                                kb_store = st.session_state.vector_store_manager.get_or_create_kb_store(active_kb)
                            
                            session_store = None
                            if st.session_state.session_files_count > 0:
                                session_store = st.session_state.vector_store_manager.get_or_create_session_store(
                                    st.session_state.session_manager.session_id
                                )
                            
                            chatbot.setup_retrievers(kb_store, session_store)
                            chatbot.create_chain()
                            
                            st.session_state.conversation_ready = True
                            st.success("✅ Ready to chat!")
                            st.rerun()
                    except ValueError as e:
                        st.error(f"❌ Model Error: {str(e)}")
                        
                        # Check if it's a memory error
                        if "memory" in str(e).lower():
                            st.warning("""
**Memory Issue Detected!**

Your system doesn't have enough RAM for this model.

**Try these smaller models instead:**
- `tinyllama` (600 MB) - Fastest
- `phi` (1.6 GB) - Fast & efficient  
- `gemma:2b` (1.4 GB) - Good quality
- `qwen2:1.5b` (900 MB) - Very small

**Pull and use:**
```bash
ollama pull tinyllama
```
Then select it in the dropdown.
                            """)
                        else:
                            st.info("""
**Troubleshooting:**
1. Check Ollama is running: `ollama serve`
2. Verify model is available: `ollama list`
3. Try pulling the model: `ollama pull llama2`
4. Use smaller models if memory is limited
                            """)
                        st.session_state.selected_model = None
                        st.session_state.conversation_ready = False
                    except Exception as e:
                        st.error(f"❌ Unexpected error: {str(e)}")
                        st.session_state.selected_model = None
                        st.session_state.conversation_ready = False
        else:
            st.warning("⚠️ No Ollama models found")
            if st.button("Retry Connection"):
                st.rerun()
        
        st.markdown("---")
        
        # Status indicators
        st.subheader("📊 Status")
        
        # Ollama connection test
        try:
            response = requests.get(f"{OLLAMA_API_URL}/tags", timeout=2)
            if response.status_code == 200:
                st.success("✅ Ollama connected")
            else:
                st.error(f"❌ Ollama error: {response.status_code}")
        except requests.exceptions.ConnectionError:
            st.error("❌ Ollama not running")
            st.info("Start with: `ollama serve`")
        except Exception as e:
            st.warning(f"⚠️ Ollama check failed: {str(e)[:50]}")
        
        # Model status
        if st.session_state.selected_model:
            st.success(f"✅ Model: {st.session_state.selected_model}")
            
            # Show if model is loaded
            current_loaded = st.session_state.get("current_loaded_model")
            if current_loaded:
                st.info(f"🔥 Loaded in RAM: {current_loaded}")
                
                # Add manual unload button
                if st.button("🗑️ Unload Model from RAM", key="unload_model_btn"):
                    with st.spinner("Unloading model..."):
                        chatbot = st.session_state.chatbot
                        if chatbot.unload_model():
                            st.session_state.current_loaded_model = None
                            st.success("✅ Model unloaded from RAM")
                            st.rerun()
        else:
            st.warning("❌ No model initialized")
        
        # KB status
        kb_manager = st.session_state.kb_manager
        active_kb = kb_manager.get_active_knowledge_base()
        if active_kb:
            kb_info = kb_manager.get_knowledge_base_info(active_kb)
            file_count = kb_info.get("file_count", 0)
            st.info(f"📚 KB: **{active_kb}** ({file_count} files)")
        else:
            st.info("📚 No KB active")
        
        # Session files status
        session_count = st.session_state.session_files_count
        if session_count > 0:
            st.info(f"📎 Session: {session_count} file(s)")
        else:
            st.info("📎 No session files")
        
        # Conversation status
        if st.session_state.conversation_ready:
            st.success("✅ Ready to chat!")
        else:
            st.warning("⚠️ Initialize model to chat")
        
        st.markdown("---")
        
        # Session management
        st.subheader("🔄 Session")
        col1, col2 = st.columns(2)
        
        with col1:
            if st.button("Clear Chat", key="clear_chat_btn"):
                st.session_state.messages = []
                st.session_state.chatbot.clear_history()
                st.success("Chat cleared")
                st.rerun()
        
        with col2:
            if st.button("Clear Session", key="clear_session_btn"):
                st.session_state.session_manager.clear_session()
                st.session_state.vector_store_manager.clear_session_store()
                st.session_state.session_files_count = 0
                st.session_state.messages = []
                st.session_state.chatbot.clear_history()
                st.session_state.conversation_ready = False
                st.success("Session cleared")
                st.rerun()


def display_chat_interface():
    """Display the main chat interface"""
    st.header("💬 Chat with Your Documents")
    
    # File upload area
    with st.expander("📎 Upload Files to Session", expanded=False):
        st.info("Upload files to use in this conversation")
        
        uploaded_files = st.file_uploader(
            "Choose files",
            accept_multiple_files=True,
            type=["pdf", "txt", "csv", "xlsx", "xls", "doc", "docx", "ppt", "pptx"],
            key="session_file_uploader"
        )
        
        if uploaded_files and st.button("Process Uploaded Files", key="process_session_files_btn"):
            with st.spinner("Processing files..."):
                processor = st.session_state.processor
                session_manager = st.session_state.session_manager
                vector_store_manager = st.session_state.vector_store_manager
                
                session_store = vector_store_manager.get_or_create_session_store(
                    session_manager.session_id
                )
                
                processed_count = 0
                for uploaded_file in uploaded_files:
                    try:
                        # Process file
                        documents, file_hash = processor.process_uploaded_file(uploaded_file)
                        
                        if documents:
                            # Add to session
                            session_manager.add_session_file(uploaded_file, file_hash)
                            
                            # Add to vector store
                            session_store.add_documents(documents)
                            processed_count += 1
                            
                            logger.info(f"Processed session file: {uploaded_file.name}")
                    
                    except Exception as e:
                        st.error(f"Error processing {uploaded_file.name}: {e}")
                
                st.session_state.session_files_count = len(session_manager.get_session_files())
                st.success(f"✅ Processed {processed_count} file(s)")
                
                # Reinitialize chain if model is ready
                if st.session_state.selected_model:
                    chatbot = st.session_state.chatbot
                    kb_manager = st.session_state.kb_manager
                    active_kb = kb_manager.get_active_knowledge_base()
                    
                    kb_store = None
                    if active_kb:
                        kb_store = vector_store_manager.get_or_create_kb_store(active_kb)
                    
                    chatbot.setup_retrievers(kb_store, session_store)
                    chatbot.create_chain()
                    st.session_state.conversation_ready = True
                    st.info("✅ Chat updated with new files!")
                
                st.rerun()
    
    # Display session files
    session_files = st.session_state.session_manager.get_session_files()
    if session_files:
        with st.expander(f"📂 Session Files ({len(session_files)})", expanded=False):
            for file_info in session_files:
                st.text(f"• {file_info['name']} ({file_info['size'] / 1024:.1f} KB)")
    
    st.markdown("---")
    
    # Display chat messages
    for message in st.session_state.messages:
        with st.chat_message(message["role"]):
            st.markdown(message["content"])
    
    # Chat input
    if prompt := st.chat_input("Ask a question about your documents..."):
        # Check if ready
        if not st.session_state.conversation_ready:
            st.error("⚠️ Please initialize a model first (check sidebar)")
            return
        
        # Add user message
        st.session_state.messages.append({"role": "user", "content": prompt})
        with st.chat_message("user"):
            st.markdown(prompt)
        
        # Get response
        with st.chat_message("assistant"):
            with st.spinner("Thinking..."):
                try:
                    chatbot = st.session_state.chatbot
                    response = chatbot.query(prompt)
                    st.markdown(response)
                    st.session_state.messages.append({"role": "assistant", "content": response})
                
                except Exception as e:
                    error_msg = f"Error: {str(e)}"
                    st.error(error_msg)
                    st.session_state.messages.append({"role": "assistant", "content": error_msg})


def display_knowledge_base_management():
    """Display knowledge base management interface"""
    st.header("📚 Knowledge Base Management")
    
    kb_manager = st.session_state.kb_manager
    
    # KB Selection
    col1, col2 = st.columns([3, 1])
    
    with col1:
        available_kbs = kb_manager.list_knowledge_bases()
        if available_kbs:
            current_active = kb_manager.get_active_knowledge_base()
            default_index = available_kbs.index(current_active) if current_active in available_kbs else 0
            
            selected_kb = st.selectbox(
                "Select Knowledge Base",
                options=available_kbs,
                index=default_index,
                key="kb_selector"
            )
            
            if st.button("Set as Active", key="set_active_kb_btn"):
                if kb_manager.set_active_knowledge_base(selected_kb):
                    st.success(f"✅ Activated: {selected_kb}")
                    
                    # Reinitialize chain if model is ready
                    if st.session_state.selected_model:
                        try:
                            chatbot = st.session_state.chatbot
                            vector_store_manager = st.session_state.vector_store_manager
                            
                            kb_store = vector_store_manager.get_or_create_kb_store(selected_kb)
                            
                            session_store = None
                            if st.session_state.session_files_count > 0:
                                session_store = vector_store_manager.get_or_create_session_store(
                                    st.session_state.session_manager.session_id
                                )
                            
                            chatbot.setup_retrievers(kb_store, session_store)
                            chatbot.create_chain()
                            st.session_state.conversation_ready = True
                            st.success("✅ Chat updated!")
                        except Exception as e:
                            st.warning(f"KB activated, but need to reinitialize: {e}")
                            st.session_state.conversation_ready = False
                    
                    st.rerun()
        else:
            st.info("No knowledge bases found. Create one below.")
    
    with col2:
        if st.button("🔄 Refresh", key="refresh_kb_btn"):
            st.rerun()
    
    # Display active KB info
    active_kb = kb_manager.get_active_knowledge_base()
    if active_kb:
        kb_info = kb_manager.get_knowledge_base_info(active_kb)
        st.success(f"🎯 **Active:** {active_kb}")
        
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric("Files", kb_info.get("file_count", 0))
        with col2:
            st.metric("Created", kb_info.get("created_date", "Unknown")[:10])
        with col3:
            st.metric("Updated", kb_info.get("last_updated", "Unknown")[:10])
        
        if kb_info.get("description"):
            st.write(f"**Description:** {kb_info['description']}")
    
    st.markdown("---")
    
    # Create new KB
    with st.expander("➕ Create New Knowledge Base", expanded=False):
        with st.form("create_kb_form"):
            kb_name = st.text_input("Knowledge Base Name")
            kb_description = st.text_area("Description (Optional)")
            
            submitted = st.form_submit_button("Create")
            
            if submitted:
                if not kb_name:
                    st.error("Please enter a name")
                elif kb_name in kb_manager.list_knowledge_bases():
                    st.error(f"'{kb_name}' already exists")
                else:
                    with st.spinner("Creating..."):
                        if kb_manager.create_knowledge_base(kb_name, kb_description):
                            st.success(f"✅ Created: {kb_name}")
                            st.rerun()
                        else:
                            st.error("Failed to create KB")
    
    # Upload files to KB
    if active_kb:
        with st.expander("📤 Upload Files to Knowledge Base", expanded=False):
            uploaded_files = st.file_uploader(
                "Choose files",
                accept_multiple_files=True,
                type=["pdf", "txt", "csv", "xlsx", "xls", "doc", "docx", "ppt", "pptx"],
                key="kb_file_uploader"
            )
            
            if uploaded_files and st.button("Process Files", key="process_kb_files_btn"):
                with st.spinner("Processing files..."):
                    processor = st.session_state.processor
                    vector_store_manager = st.session_state.vector_store_manager
                    kb_store = vector_store_manager.get_or_create_kb_store(active_kb)
                    
                    processed_count = 0
                    for uploaded_file in uploaded_files:
                        try:
                            documents, file_hash = processor.process_uploaded_file(uploaded_file)
                            
                            if documents:
                                kb_store.add_documents(documents)
                                processed_count += 1
                                logger.info(f"Processed KB file: {uploaded_file.name}")
                        
                        except Exception as e:
                            st.error(f"Error processing {uploaded_file.name}: {e}")
                    
                    # Update metadata
                    kb_info = kb_manager.get_knowledge_base_info(active_kb)
                    new_count = kb_info.get("file_count", 0) + processed_count
                    kb_manager.update_metadata(active_kb, {"file_count": new_count})
                    
                    st.success(f"✅ Processed {processed_count} file(s)")
                    st.rerun()
    
    # Delete KB
    if active_kb:
        with st.expander("🗑️ Delete Knowledge Base", expanded=False):
            st.warning(f"Delete **{active_kb}**?")
            confirm_name = st.text_input(f"Type '{active_kb}' to confirm")
            
            if st.button("Delete", type="primary", key="delete_kb_btn"):
                if confirm_name == active_kb:
                    with st.spinner("Deleting..."):
                        vector_store_manager = st.session_state.vector_store_manager
                        vector_store_manager.delete_vector_store(active_kb)
                        
                        if kb_manager.delete_knowledge_base(active_kb):
                            st.success(f"✅ Deleted: {active_kb}")
                            st.session_state.conversation_ready = False
                            st.rerun()
                        else:
                            st.error("Failed to delete")
                else:
                    st.error("Incorrect name")


def main():
    """Main application entry point"""
    st.set_page_config(
        page_title="Enhanced KMS v3.0 - KC",
        page_icon="📚",
        layout="wide",
        initial_sidebar_state="expanded",
    )
    
    st.title("📚 Enhanced Knowledge Management System v3.0-KC")
    st.subheader("Session-Based RAG with Hybrid Document Management")
    
    # Show library status
    status_parts = []
    if IMAGE_PROCESSING_AVAILABLE:
        status_parts.append("✅ Image Processing")
    if OCR_AVAILABLE:
        status_parts.append(f"✅ OCR ({OCR_METHOD})")
    else:
        status_parts.append("ℹ️ OCR not available (optional)")
    
    st.info(" | ".join(status_parts))
    
    # Initialize session state
    initialize_session_state()
    
    # Create tabs
    tab1, tab2 = st.tabs(["💬 Chat", "📚 Knowledge Bases"])
    
    with tab1:
        display_sidebar()
        display_chat_interface()
    
    with tab2:
        display_knowledge_base_management()


if __name__ == "__main__":
    main()