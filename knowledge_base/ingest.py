import os
import asyncio
from typing import List
from langchain_community.document_loaders import (
    PyPDFLoader, 
    Docx2txtLoader, 
    TextLoader, 
    CSVLoader,
    UnstructuredExcelLoader
)
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_postgres.vectorstores import PGVector
from langchain_ollama import OllamaEmbeddings
from sqlalchemy import select, update
from storage.database import AsyncSessionLocal, engine
from storage.models import DocumentMetadata
import uuid
from dotenv import load_dotenv

load_dotenv()

OLLAMA_URL = os.getenv("OLLAMA_URL", "http://localhost:11434")
POSTGRES_USER = os.getenv("POSTGRES_USER", "rag_user")
POSTGRES_PASSWORD = os.getenv("POSTGRES_PASSWORD", "rag_pass")
POSTGRES_DB = os.getenv("POSTGRES_DB", "rag_db")
POSTGRES_HOST = os.getenv("POSTGRES_HOST", "localhost")
POSTGRES_PORT = os.getenv("POSTGRES_PORT", "5433")

# Check for a full CONNECTION_STRING (postgresql+psycopg://...)
# If not provided, build it from parts.
CONNECTION_STRING = os.getenv("VECTOR_DB_URL")
if not CONNECTION_STRING:
    CONNECTION_STRING = f"postgresql+psycopg://{POSTGRES_USER}:{POSTGRES_PASSWORD}@{POSTGRES_HOST}:{POSTGRES_PORT}/{POSTGRES_DB}"

COLLECTION_NAME = os.getenv("COLLECTION_NAME", "industrial_docs")

embeddings = OllamaEmbeddings(
    model=os.getenv("OLLAMA_EMBED_MODEL", "nomic-embed-text:latest"), 
    base_url=OLLAMA_URL
)

# OCR Initialization logic from expected code
IMAGE_PROCESSING_AVAILABLE = False
OCR_AVAILABLE = False
OCR_METHOD = None

try:
    from PIL import Image as PILImage
    IMAGE_PROCESSING_AVAILABLE = True
    try:
        import easyocr
        OCR_AVAILABLE = True
        OCR_METHOD = "easyocr"
    except ImportError:
        try:
            import pytesseract
            # Check if tesseract is installed
            pytesseract.get_tesseract_version()
            OCR_AVAILABLE = True
            OCR_METHOD = "tesseract"
        except Exception:
            pass
except ImportError:
    pass

class OCRProcessor:
    """Optimized OCR processor from the expected implementation"""
    def __init__(self):
        self.ocr_reader = None
        self.ocr_method = OCR_METHOD
        self._initialize_ocr()
    
    def _initialize_ocr(self):
        if not OCR_AVAILABLE:
            return
        try:
            if self.ocr_method == "easyocr":
                import easyocr
                self.ocr_reader = easyocr.Reader(["en"], gpu=False)
            elif self.ocr_method == "tesseract":
                import pytesseract
                pass # Already checked availability
        except Exception as e:
            print(f"Error initializing OCR: {e}")
            self.ocr_reader = None
    
    def extract_text_from_image(self, image_path: str) -> str:
        if not OCR_AVAILABLE or not IMAGE_PROCESSING_AVAILABLE:
            return ""
        try:
            image = PILImage.open(image_path)
            if image.mode != "RGB":
                image = image.convert("RGB")
            
            text = ""
            if self.ocr_method == "easyocr" and self.ocr_reader:
                results = self.ocr_reader.readtext(image_path)
                text_parts = [text_content for bbox, text_content, confidence in results if confidence > 0.5]
                text = " ".join(text_parts)
            elif self.ocr_method == "tesseract":
                import pytesseract
                text = pytesseract.image_to_string(image)
            return text.strip()
        except Exception as e:
            print(f"Error in OCR processing: {e}")
            return ""

ocr_processor = OCRProcessor()

class ExcelLoader:
    """Excel loader from the expected implementation"""
    def __init__(self, file_path):
        self.file_path = file_path
    def load(self):
        try:
            import openpyxl
            workbook = openpyxl.load_workbook(self.path if hasattr(self, 'path') else self.file_path, data_only=True)
            documents = []
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                content_parts = [f"Sheet: {sheet_name}\n"]
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        content_parts.append(row_text)
                content = "\n".join(content_parts)
                if content.strip():
                    doc = Document(page_content=content, metadata={"source": os.path.basename(self.file_path), "sheet": sheet_name})
                    documents.append(doc)
            return documents
        except Exception as e:
            print(f"Error loading Excel file {self.file_path}: {e}")
            return []

class WordLoader:
    """Word loader from the expected implementation"""
    def __init__(self, file_path):
        self.file_path = file_path
    def load(self):
        try:
            from docx import Document as DocxDocument
            doc = DocxDocument(self.file_path)
            content_parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    content_parts.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    if row_text.strip():
                        content_parts.append(row_text)
            content = "\n\n".join(content_parts)
            if content.strip():
                return [Document(page_content=content, metadata={"source": os.path.basename(self.file_path)})]
            return []
        except Exception as e:
            print(f"Error loading Word file {self.file_path}: {e}")
            return []

class PowerPointLoader:
    """PowerPoint loader from the expected implementation"""
    def __init__(self, file_path):
        self.file_path = file_path
    def load(self):
        try:
            from pptx import Presentation
            prs = Presentation(self.file_path)
            documents = []
            for slide_num, slide in enumerate(prs.slides, 1):
                content_parts = [f"Slide {slide_num}:"]
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text)
                content = "\n".join(content_parts)
                if content.strip():
                    doc = Document(page_content=content, metadata={"source": os.path.basename(self.file_path), "slide": slide_num})
                    documents.append(doc)
            return documents
        except Exception as e:
            print(f"Error loading PowerPoint file {self.file_path}: {e}")
            return []

def get_loader(file_path: str):
    ext = os.path.splitext(file_path)[1].lower()
    print(f"Selecting loader for extension: {ext} ({file_path})")
    if ext == ".pdf":
        return PyPDFLoader(file_path)
    elif ext in [".docx", ".doc"]:
        return WordLoader(file_path)
    elif ext in [".xlsx", ".xls"]:
        return ExcelLoader(file_path)
    elif ext == ".csv":
        return CSVLoader(file_path)
    elif ext in [".pptx", ".ppt"]:
        return PowerPointLoader(file_path)
    elif ext in [".txt", ".md"]:
        class RobustTextLoader:
            def __init__(self, path):
                self.path = path
            def load(self):
                encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252"]
                for enc in encodings:
                    try:
                        loader = TextLoader(self.path, encoding=enc)
                        return loader.load()
                    except UnicodeDecodeError:
                        continue
                raise ValueError(f"Could not decode {self.path} with common encodings.")
        return RobustTextLoader(file_path)
    elif ext in [".jpg", ".jpeg", ".png", ".bmp"]:
        class ImageLoader:
            def __init__(self, path):
                self.path = path
            def load(self):
                text = ocr_processor.extract_text_from_image(self.path)
                if text:
                    return [Document(page_content=text, metadata={"source": os.path.basename(self.path)})]
                return []
        return ImageLoader(file_path)
    else:
        raise ValueError(f"Unsupported file extension: {ext}")

async def process_document(pdf_file: str, file_path: str, session_id: str = None):
    """Processes a single document: splits text, embeds, and saves to PGVector."""
    print(f"Starting process_document for: {pdf_file} (Session: {session_id})")
    async with AsyncSessionLocal() as session:
        # Check if already processed
        stmt = select(DocumentMetadata).where(DocumentMetadata.filename == pdf_file)
        result = await session.execute(stmt)
        meta = result.scalar_one_or_none()
        
        if meta and meta.status == "processed":
            # If reprocessing with different session, might want to update, but skipping for now or just log
            print(f"Skipping {pdf_file}, already processed.")
            return

        if not meta:
            print(f"Creating new metadata for {pdf_file}")
            meta = DocumentMetadata(
                filename=pdf_file,
                file_path=file_path,
                status="processing",
                session_id=session_id
            )
            session.add(meta)
            await session.commit()
        
        print(f"Processing {pdf_file}...")
        
        try:
            loader = get_loader(file_path)
            print(f"Loading document with {loader.__class__.__name__}...")
            docs = await asyncio.to_thread(loader.load)
            print(f"Loaded {len(docs)} document objects.")
            
            text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
            splits = await asyncio.to_thread(text_splitter.split_documents, docs)
            
            # Ensure consistent metadata for deletion and search
            for split in splits:
                split.metadata["source"] = pdf_file
                if session_id:
                    split.metadata["session_id"] = session_id
                # Prepend source filename to content to aid retrieval
                split.page_content = f"--- Document: {pdf_file} ---\n{split.page_content}"
                
            print(f"Created {len(splits)} splits for {pdf_file}")
            
            def add_to_vectorstore():
                print(f"Connecting to vector store for {pdf_file}...")
                vector_store = PGVector(
                    embeddings=embeddings,
                    collection_name=COLLECTION_NAME,
                    connection=CONNECTION_STRING,
                    use_jsonb=True,
                )
                
                # Add to vector store in batches to avoid overwhelming the system
                batch_size = 50
                for i in range(0, len(splits), batch_size):
                    batch = splits[i:i+batch_size]
                    print(f"Adding batch {i//batch_size + 1}/{(len(splits)-1)//batch_size + 1} ({len(batch)} splits) for {pdf_file}...")
                    vector_store.add_documents(batch)
                
            print(f"Starting vector store ingestion for {pdf_file}...")
            await asyncio.to_thread(add_to_vectorstore)
            print(f"Successfully added {pdf_file} to vector store.")
            
            # Update status
            meta.status = "processed"
            await session.merge(meta)
            await session.commit()
            print(f"Successfully updated status for {pdf_file}")
            
        except Exception as e:
            print(f"!!! Error processing {pdf_file}: {e}")
            import traceback
            traceback.print_exc()
            meta.status = "error"
            await session.merge(meta)
            await session.commit()

async def ingest_pdfs():
    doc_dir = "knowledge_base/documents"
    if not os.path.exists(doc_dir):
        print(f"Directory {doc_dir} not found.")
        return

    SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".doc", ".txt", ".md", ".csv", ".xlsx", ".xls", ".jpg", ".jpeg", ".png", ".bmp", ".pptx", ".ppt"}
    pdf_files = [f for f in os.listdir(doc_dir) if os.path.splitext(f)[1].lower() in SUPPORTED_EXTENSIONS]
    
    for pdf_file in pdf_files:
        file_path = os.path.join(doc_dir, pdf_file)
        await process_document(pdf_file, file_path)

if __name__ == "__main__":
    asyncio.run(ingest_pdfs())
